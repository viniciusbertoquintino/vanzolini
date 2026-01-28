from pptx import Presentation

prs = Presentation("template.pptx")

for i, slide in enumerate(prs.slides, 1):
    print(f"\n=== SLIDE {i} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            print("Texto do shape:")
            print(repr(shape.text))
