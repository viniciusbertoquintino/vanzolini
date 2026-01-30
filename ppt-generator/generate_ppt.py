from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

LAYOUTS = {
    "CAPA": 0,
    "1COL": 1,
    "2COL": 2,
    "IMG": 3
}

def generate_ppt(template_path, content_path, output_path):
    prs = Presentation(template_path)

    with open(content_path, "r", encoding="utf-8") as f:
        raw = f.read()

    sections = [s.strip() for s in raw.split("#") if s.strip()]

    for section in sections:
        lines = [l.strip() for l in section.splitlines() if l.strip()]

        tipo = lines[0]
        dados = lines[1:]

        layout = LAYOUTS.get(tipo)
        if layout is None:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[layout])

        # CAPA
        if tipo == "CAPA":
            slide.shapes.title.text = dados[0]
            slide.placeholders[1].text = dados[1]

        # 1 COLUNA
        elif tipo == "1COL":
            slide.shapes.title.text = dados[0]
            slide.placeholders[1].text = dados[1]

        # 2 COLUNAS
        elif tipo == "2COL":
            slide.shapes.title.text = dados[0]
            slide.placeholders[1].text = dados[1]
            slide.placeholders[2].text = dados[2]

        # SLIDE COM IMAGEM
        elif tipo == "IMG":
            slide.shapes.title.text = dados[0]
            img_path = dados[1]
            slide.shapes.add_picture(img_path, Inches(1), Inches(2))
            slide.placeholders[1].text = dados[2]

    prs.save(output_path)
